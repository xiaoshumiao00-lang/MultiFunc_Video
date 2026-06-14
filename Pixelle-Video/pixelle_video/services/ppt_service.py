# Copyright (C) 2025 AIDC-AI
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#     http://www.apache.org/licenses/LICENSE-2.0
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""
PPT/PDF Teaching Material Parsing Service

Supports:
- .pptx file parsing and slide rendering (via python-pptx + Pillow fallback)
- .pdf file parsing and page rendering (via PyMuPDF)
- Unified output: list of slides with image path and extracted text
"""

import os
import re
import io
import uuid
from pathlib import Path
from dataclasses import dataclass, field
from typing import List, Optional, Tuple, Literal

from loguru import logger
from PIL import Image, ImageDraw, ImageFont

from pixelle_video.utils.os_util import get_temp_path, ensure_dir


@dataclass
class SlideInfo:
    """Represents a single slide/page of teaching material"""
    index: int
    image_path: str
    text: str
    title: str = ""
    elements: List[dict] = field(default_factory=list)


class PPTService:
    """
    Service for parsing and rendering PPT/PDF teaching materials.
    
    For PDF: uses PyMuPDF for high-fidelity page rendering.
    For PPTX: uses python-pptx to extract content and Pillow to render a
    simplified slide representation (since Office/LibreOffice is not available).
    """
    
    # Standard slide rendering resolution (16:9)
    SLIDE_WIDTH = 1920
    SLIDE_HEIGHT = 1080
    
    # Fallback fonts for different platforms
    FONT_CANDIDATES = [
        # Windows
        "C:/Windows/Fonts/msyh.ttc",
        "C:/Windows/Fonts/simhei.ttf",
        "C:/Windows/Fonts/simsun.ttc",
        # Linux
        "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        # macOS
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/STHeiti Light.ttc",
    ]
    
    def __init__(self):
        self._font = None
        self._title_font = None
        self._body_font = None
    
    def parse(self, file_path: str, output_dir: Optional[str] = None) -> List[SlideInfo]:
        """
        Parse a PPT/PDF file and return a list of SlideInfo.
        
        Args:
            file_path: Path to .pptx or .pdf file
            output_dir: Directory to save rendered slide images (default: temp/ppt_slides)
        
        Returns:
            List of SlideInfo objects
        """
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"Teaching material file not found: {file_path}")
        
        suffix = file_path.suffix.lower()
        if suffix == ".pdf":
            return self._parse_pdf(str(file_path), output_dir)
        elif suffix in (".pptx", ".ppt"):
            return self._parse_pptx(str(file_path), output_dir)
        else:
            raise ValueError(f"Unsupported file format: {suffix}. Only .pptx and .pdf are supported.")
    
    def _get_output_dir(self, output_dir: Optional[str] = None) -> str:
        """Get or create output directory for slide images"""
        if output_dir:
            return ensure_dir(output_dir)
        unique_id = uuid.uuid4().hex[:8]
        return ensure_dir(get_temp_path(f"ppt_slides_{unique_id}"))
    
    # ======================================================================
    # PDF Parsing
    # ======================================================================
    def _parse_pdf(self, file_path: str, output_dir: Optional[str] = None) -> List[SlideInfo]:
        """Parse PDF file and render each page as image"""
        try:
            import fitz  # PyMuPDF
        except ImportError as e:
            raise RuntimeError("PyMuPDF (fitz) is required for PDF parsing. Please install it.") from e
        
        output_dir = self._get_output_dir(output_dir)
        slides = []
        
        logger.info(f"Parsing PDF: {file_path}")
        doc = fitz.open(file_path)
        
        try:
            for page_idx in range(len(doc)):
                page = doc.load_page(page_idx)
                
                # Render page at high resolution
                mat = fitz.Matrix(2.0, 2.0)  # 2x scale for high quality
                pix = page.get_pixmap(matrix=mat, alpha=False)
                
                # Save image
                image_path = os.path.join(output_dir, f"slide_{page_idx:03d}.png")
                pix.save(image_path)
                
                # Extract text
                text = page.get_text("text").strip()
                title = self._extract_title(text)
                
                slides.append(SlideInfo(
                    index=page_idx,
                    image_path=image_path,
                    text=text,
                    title=title
                ))
                
                logger.debug(f"PDF page {page_idx + 1}/{len(doc)} rendered: {image_path}")
        finally:
            doc.close()
        
        logger.info(f"PDF parsed: {len(slides)} slides")
        return slides
    
    # ======================================================================
    # PPTX Parsing
    # ======================================================================
    def _parse_pptx(self, file_path: str, output_dir: Optional[str] = None) -> List[SlideInfo]:
        """
        Parse PPTX file and render each slide as image.
        
        Note: Since Office/LibreOffice is not available, this uses python-pptx
        to extract content and Pillow to render a simplified representation.
        For best visual fidelity, users are recommended to export PPT as PDF first.
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError as e:
            raise RuntimeError("python-pptx is required for PPTX parsing. Please install it.") from e
        
        output_dir = self._get_output_dir(output_dir)
        slides = []
        
        logger.info(f"Parsing PPTX: {file_path}")
        prs = Presentation(file_path)
        
        for slide_idx, slide in enumerate(prs.slides):
            # Extract text and elements
            text_parts = []
            elements = []
            title = ""
            
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                shape_text = shape.text_frame.text.strip()
                if not shape_text:
                    continue
                
                text_parts.append(shape_text)
                
                # Identify title (usually first text shape or placeholder)
                if shape.is_placeholder:
                    from pptx.enum.shapes import PP_PLACEHOLDER
                    if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                        title = shape_text
                
                elements.append({
                    "text": shape_text,
                    "is_placeholder": shape.is_placeholder,
                    "left": shape.left.inches if hasattr(shape, "left") else 0,
                    "top": shape.top.inches if hasattr(shape, "top") else 0,
                    "width": shape.width.inches if hasattr(shape, "width") else 1,
                    "height": shape.height.inches if hasattr(shape, "height") else 0.5,
                })
            
            # Extract images from shapes
            images = self._extract_pptx_images(slide, output_dir, slide_idx)
            
            # Render slide image
            full_text = "\n".join(text_parts)
            if not title and full_text:
                title = self._extract_title(full_text)
            
            image_path = self._render_pptx_slide(
                slide_idx=slide_idx,
                title=title,
                text_parts=text_parts,
                elements=elements,
                images=images,
                output_dir=output_dir
            )
            
            slides.append(SlideInfo(
                index=slide_idx,
                image_path=image_path,
                text=full_text,
                title=title,
                elements=elements
            ))
            
            logger.debug(f"PPTX slide {slide_idx + 1}/{len(prs.slides)} rendered: {image_path}")
        
        logger.info(f"PPTX parsed: {len(slides)} slides")
        return slides
    
    def _extract_pptx_images(
        self,
        slide,
        output_dir: str,
        slide_idx: int
    ) -> List[Tuple[str, float, float, float, float]]:
        """Extract embedded images from a PPTX slide"""
        images = []
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            return images
        
        img_counter = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    image_bytes = image.blob
                    ext = image.ext
                    if not ext:
                        ext = "png"
                    
                    img_filename = f"slide_{slide_idx:03d}_img_{img_counter}.{ext}"
                    img_path = os.path.join(output_dir, img_filename)
                    with open(img_path, "wb") as f:
                        f.write(image_bytes)
                    
                    images.append((
                        img_path,
                        shape.left.inches if hasattr(shape, "left") else 0,
                        shape.top.inches if hasattr(shape, "top") else 0,
                        shape.width.inches if hasattr(shape, "width") else 1,
                        shape.height.inches if hasattr(shape, "height") else 1,
                    ))
                    img_counter += 1
                except Exception as e:
                    logger.warning(f"Failed to extract image from PPTX slide {slide_idx}: {e}")
        
        return images
    
    def _render_pptx_slide(
        self,
        slide_idx: int,
        title: str,
        text_parts: List[str],
        elements: List[dict],
        images: List[Tuple[str, float, float, float, float]],
        output_dir: str
    ) -> str:
        """Render a PPTX slide as an image using Pillow"""
        width, height = self.SLIDE_WIDTH, self.SLIDE_HEIGHT
        
        # Create white background
        img = Image.new("RGB", (width, height), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)
        
        # Load fonts
        title_font = self._get_font(size=54, bold=True)
        body_font = self._get_font(size=36)
        caption_font = self._get_font(size=24)
        
        # Draw subtle header bar
        draw.rectangle([0, 0, width, 120], fill=(245, 247, 250))
        
        # Draw title
        if title:
            draw.text((60, 60), title, fill=(33, 37, 41), font=title_font, anchor="lm")
        
        # Draw content area
        y_offset = 160
        max_width = width - 120
        line_height = 50
        
        # If we have structured elements with positions, try to use them
        if elements:
            # Sort by vertical position
            sorted_elements = sorted(
                [e for e in elements if e.get("text")],
                key=lambda x: x.get("top", 0)
            )
            
            for elem in sorted_elements:
                text = elem["text"]
                # Skip title if already drawn
                if text == title:
                    continue
                
                # Wrap text
                wrapped_lines = self._wrap_text(text, body_font, max_width)
                for line in wrapped_lines:
                    if y_offset + line_height > height - 60:
                        break
                    draw.text((60, y_offset), line, fill=(52, 58, 64), font=body_font)
                    y_offset += line_height
                y_offset += 20  # spacing between elements
        else:
            # Fallback: draw all text
            for text in text_parts:
                wrapped_lines = self._wrap_text(text, body_font, max_width)
                for line in wrapped_lines:
                    if y_offset + line_height > height - 60:
                        break
                    draw.text((60, y_offset), line, fill=(52, 58, 64), font=body_font)
                    y_offset += line_height
                y_offset += 20
        
        # Draw extracted images (simplified: place in right column if space allows)
        if images:
            col_x = width - 520
            col_y = 160
            for img_path, _, _, _, _ in images[:2]:  # Limit to first 2 images
                try:
                    slide_img = Image.open(img_path)
                    # Resize to fit column
                    slide_img.thumbnail((480, 400), Image.Resampling.LANCZOS)
                    if col_y + slide_img.height < height - 60:
                        img.paste(slide_img, (col_x, col_y))
                        col_y += slide_img.height + 30
                except Exception as e:
                    logger.warning(f"Failed to paste image {img_path}: {e}")
        
        # Add slide number
        draw.text((width - 40, height - 30), f"{slide_idx + 1}",
                  fill=(150, 150, 150), font=caption_font, anchor="rb")
        
        # Save
        image_path = os.path.join(output_dir, f"slide_{slide_idx:03d}.png")
        img.save(image_path, "PNG")
        return image_path
    
    # ======================================================================
    # Helpers
    # ======================================================================
    def _extract_title(self, text: str) -> str:
        """Extract title from slide text (first non-empty line)"""
        lines = [line.strip() for line in text.split("\n") if line.strip()]
        if not lines:
            return ""
        title = lines[0]
        # Limit title length
        if len(title) > 100:
            title = title[:97] + "..."
        return title
    
    def _get_font(self, size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
        """Get a usable font, falling back through candidates"""
        if self._font is None:
            for font_path in self.FONT_CANDIDATES:
                if os.path.exists(font_path):
                    try:
                        self._font = font_path
                        break
                    except Exception:
                        continue
            if self._font is None:
                # Use default Pillow font
                self._font = None
        
        if self._font:
            try:
                return ImageFont.truetype(self._font, size)
            except Exception:
                pass
        
        return ImageFont.load_default()
    
    def _wrap_text(self, text: str, font: ImageFont.FreeTypeFont, max_width: int) -> List[str]:
        """Wrap text to fit within max_width"""
        if not text:
            return []
        
        lines = []
        for paragraph in text.split("\n"):
            if not paragraph.strip():
                continue
            
            words = []
            # For CJK, split by character; for others, split by word
            if self._is_cjk(paragraph):
                words = list(paragraph)
            else:
                words = paragraph.split(" ")
            
            current_line = ""
            for word in words:
                test_line = current_line + word if not current_line else current_line + (" " if not self._is_cjk(word) else "") + word
                bbox = font.getbbox(test_line)
                text_width = bbox[2] - bbox[0] if bbox else 0
                
                if text_width <= max_width:
                    current_line = test_line
                else:
                    if current_line:
                        lines.append(current_line)
                    current_line = word
            
            if current_line:
                lines.append(current_line)
        
        return lines
    
    @staticmethod
    def _is_cjk(text: str) -> bool:
        """Check if text contains CJK characters"""
        return bool(re.search(r"[\u4e00-\u9fff]", text))
    
    def get_font_path(self) -> Optional[str]:
        """Return the resolved font path for external use"""
        if self._font is None:
            self._get_font(24)
        return self._font


def parse_teaching_material(file_path: str, output_dir: Optional[str] = None) -> List[SlideInfo]:
    """Convenience function to parse PPT/PDF teaching material"""
    service = PPTService()
    return service.parse(file_path, output_dir)
